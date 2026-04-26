from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

prs = Presentation("PPT Template.pptx")

def set_text(shape, text, font_size=None, bold=None, color=None, alignment=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)

def add_bullet_text(shape, items, font_size=10, color=None, title=None, title_size=11):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        if color:
            run.font.color.rgb = RGBColor(*color)
        first = False

    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = RGBColor(*color)

def add_text_box(slide, text, left, top, width, height, font_size=10, bold=False, color=(0,0,0), alignment=PP_ALIGN.LEFT, word_wrap=True):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return txBox

def add_multi_para_box(slide, paragraphs, left, top, width, height, font_size=10, color=(50,50,50)):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (bold, text) in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
    return txBox

# ─────────────────────────────────────────────
# SLIDE 1 – Title / Cover
# ─────────────────────────────────────────────
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        if "Team Name" in shape.text_frame.text:
            set_text(shape, "YellowSense Technologies", font_size=22, bold=True, color=(11, 83, 148))

# Add team details as text boxes on slide 1
# Project name
add_text_box(slide1, "SUVIDHA Connect",
    left=2639877, top=3500000, width=4500000, height=350000,
    font_size=18, bold=True, color=(11, 83, 148), alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide1, "Smart Unified Virtual Interface for Digital Helpdesk & Assistance",
    left=1800000, top=3850000, width=5800000, height=280000,
    font_size=11, bold=False, color=(80, 80, 80), alignment=PP_ALIGN.CENTER)

# Team info left column
add_text_box(slide1, "Team ID: NexaCore | Slot 39 | 05:20 PM – 05:30 PM",
    left=527050, top=4016225, width=4200000, height=280000,
    font_size=10, bold=False, color=(60, 60, 60))

# Demo links
add_multi_para_box(slide1, [
    (True,  "🌐 Live Demo:"),
    (False, "https://suvidha-connect.vercel.app"),
    (True,  "💻 GitHub:"),
    (False, "github.com/YellowSense/suvidha-connect"),
], left=527050, top=4300000, width=4200000, height=500000, font_size=9, color=(30, 30, 30))

# ─────────────────────────────────────────────
# SLIDE 2 – Problem Understanding
# ─────────────────────────────────────────────
slide2 = prs.slides[1]

# Main content area – two columns
# Left column: Problem Statement
add_multi_para_box(slide2, [
    (True,  "🔴 The Civic Service Gap"),
    (False, ""),
    (False, "• 70%+ citizens in Tier-2/3 cities lack digital access to utility services"),
    (False, "• Long queues at electricity, gas & municipal offices waste 2–4 hrs/visit"),
    (False, "• No unified platform — citizens visit 3+ offices for basic services"),
    (False, "• Language barriers exclude non-English speakers from digital portals"),
    (False, "• Elderly & differently-abled citizens have no accessible kiosk option"),
    (False, ""),
    (True,  "🎯 Target Services"),
    (False, ""),
    (False, "• Electricity: Bill payment, new connections, consumption analytics"),
    (False, "• Gas (Assam Gas / LPG): Refill booking, subsidy tracking, safety alerts"),
    (False, "• Municipal: Grievance filing, waste management, document services"),
    (False, ""),
    (True,  "📍 Target Users"),
    (False, ""),
    (False, "• Semi-literate & first-time digital users"),
    (False, "• Senior citizens & Divyangjan (differently-abled)"),
    (False, "• Rural & peri-urban citizens in Assam & NE India"),
], left=311699, top=700000, width=4400000, height=3900000, font_size=10, color=(30, 30, 30))

# Right column: Impact stats
add_multi_para_box(slide2, [
    (True,  "📊 Problem Scale"),
    (False, ""),
    (False, "• 500M+ citizens rely on govt utility services"),
    (False, "• ₹2,000 Cr+ lost annually in inefficient billing"),
    (False, "• 60% complaints go unresolved within SLA"),
    (False, "• Only 12% rural citizens use digital payment"),
    (False, ""),
    (True,  "⚠️ Key Pain Points"),
    (False, ""),
    (False, "• No single-window service delivery"),
    (False, "• Lack of real-time complaint tracking"),
    (False, "• No offline-capable kiosk solution"),
    (False, "• Zero gamification for digital adoption"),
    (False, "• No AI-assisted navigation for illiterate users"),
    (False, ""),
    (True,  "✅ SUVIDHA Connect Addresses All of These"),
], left=4900000, top=700000, width=4000000, height=3900000, font_size=10, color=(30, 30, 30))

# ─────────────────────────────────────────────
# SLIDE 3 – Proposed Solution
# ─────────────────────────────────────────────
slide3 = prs.slides[2]

add_multi_para_box(slide3, [
    (True,  "✅ Prototype Maturity: Fully Functional System (MVP+)"),
    (False, ""),
    (False, "Live at: https://suvidha-connect.vercel.app"),
    (False, ""),
    (True,  "⚡ Electricity Features"),
    (False, "• Bill fetch & payment via Consumer ID"),
    (False, "• Consumption analytics with monthly comparison"),
    (False, "• New connection request with document upload"),
    (False, "• Tamper-proof digital receipt with QR code"),
    (False, ""),
    (True,  "🔥 Gas (Assam Gas / LPG) Features"),
    (False, "• One-touch refill booking & payment"),
    (False, "• DBTL subsidy status tracking"),
    (False, "• Real-time safety alerts & quick leak report"),
    (False, "• Piped gas connection management"),
    (False, ""),
    (True,  "🏛️ Municipal / Civic Features"),
    (False, "• Smart complaint filing with geo-tag & tracking ID"),
    (False, "• Waste management schedule & fee payment"),
    (False, "• Birth/Death/Caste certificate apply & download"),
    (False, "• Real-time civic alerts broadcast ticker"),
], left=311699, top=650000, width=4400000, height=4100000, font_size=10, color=(30, 30, 30))

add_multi_para_box(slide3, [
    (True,  "🔐 Authentication"),
    (False, "• Mobile OTP login (9876543210 / OTP: 123456)"),
    (False, "• Consumer ID + PIN login"),
    (False, "• QR code scan (placeholder ready)"),
    (False, ""),
    (True,  "🤖 AI & Smart Features"),
    (False, "• SUVIDHA Sahayak – NLP chat assistant"),
    (False, "• Voice Commander (Web Speech API)"),
    (False, "• Handles Hindi queries: 'Mera bijli bill dikhao'"),
    (False, ""),
    (True,  "🏆 Gamification"),
    (False, "• SUVIDHA Points for on-time bill payment"),
    (False, "• Redeemable for transport passes & tax rebates"),
    (False, ""),
    (True,  "📅 Additional Services"),
    (False, "• Appointment booking (reduces office queues)"),
    (False, "• Rewards dashboard & leaderboard"),
    (False, "• Admin monitoring dashboard with analytics"),
    (False, "• Audit logs & integrity ledger (blockchain-lite)"),
], left=4900000, top=650000, width=4000000, height=4100000, font_size=10, color=(30, 30, 30))

# ─────────────────────────────────────────────
# SLIDE 4 – System Architecture
# ─────────────────────────────────────────────
slide4 = prs.slides[3]

add_multi_para_box(slide4, [
    (True,  "🏗️ 3-Layer Architecture"),
    (False, ""),
    (True,  "① Presentation Layer"),
    (False, "  React 18 + TypeScript + TailwindCSS + Shadcn UI"),
    (False, "  Modules: Login | Services | Admin | Voice | Chat | VirtualKeyboard"),
    (False, ""),
    (True,  "② Business Logic Layer"),
    (False, "  React Context API (Auth, Kiosk, Keyboard, i18n Contexts)"),
    (False, "  Custom hooks for session management & accessibility"),
    (False, ""),
    (True,  "③ Data Layer"),
    (False, "  TypeScript-typed mock services (PostgreSQL-schema-ready)"),
    (False, "  Bill Service | Complaint Service | Request Service | Alert Service"),
    (False, ""),
    (True,  "☁️ Deployment Strategy"),
    (False, "  • Cloud: Vercel (current) — 99.9% uptime, global CDN"),
    (False, "  • On-Premise: Deployable on NIC servers / state data centres"),
    (False, "  • Hybrid: Static frontend + REST API backend (future)"),
    (False, ""),
    (True,  "📡 Internet Dependency: Low–Medium"),
    (False, "  PWA caches critical UI — works offline for bill viewing & forms"),
], left=311699, top=650000, width=4400000, height=4100000, font_size=10, color=(30, 30, 30))

add_multi_para_box(slide4, [
    (True,  "🖥️ Infrastructure Requirements"),
    (False, ""),
    (False, "• Kiosk Hardware: Android tablet / Raspberry Pi 4 / x86 mini-PC"),
    (False, "• RAM: 2GB minimum (4GB recommended)"),
    (False, "• Storage: 512MB for PWA cache"),
    (False, "• Display: 10–15\" capacitive touchscreen"),
    (False, "• Network: 4G/LTE or broadband (fallback: offline PWA)"),
    (False, ""),
    (True,  "🔄 CI/CD & Scalability"),
    (False, ""),
    (False, "• Git-based CI/CD via Vercel — push-to-deploy"),
    (False, "• Updates propagate to all kiosks simultaneously"),
    (False, "• Modular React components — add departments without rebuild"),
    (False, "• Code splitting + lazy loading for fast kiosk boot"),
    (False, ""),
    (True,  "🔌 Future Backend Integrations"),
    (False, ""),
    (False, "• REST APIs → State DISCOM / Gas utility databases"),
    (False, "• Payment Gateway → UPI / Bharat BillPay"),
    (False, "• PostgreSQL → Persistent citizen data"),
    (False, "• SMS/Email → OTP & receipt delivery"),
], left=4900000, top=650000, width=4000000, height=4100000, font_size=10, color=(30, 30, 30))

# ─────────────────────────────────────────────
# SLIDE 5 – UI/UX, Security, Scalability & Compliance
# ─────────────────────────────────────────────
slide5 = prs.slides[4]

# Use the existing Subtitle 2 placeholder for main content
for shape in slide5.shapes:
    if shape.name == "Subtitle 2":
        add_bullet_text(shape, [
            "🖐️ TOUCH-OPTIMISED UI",
            "  • Min 56px touch targets — works on resistive & capacitive screens",
            "  • Custom VirtualKeyboard.tsx — no OS keyboard dependency",
            "  • One-Task-Per-Screen flow: Identify → Verify → Pay → Receipt",
            "  • High-contrast Govt theme: Blue / Saffron / Green",
            "",
            "♿ ACCESSIBILITY & INCLUSION",
            "  • Bilingual: English + Hindi (i18next) — regional languages ready",
            "  • Text-to-Speech (TTS) reads instructions aloud",
            "  • Voice Commander — fully hands-free kiosk navigation",
            "  • High-Contrast Mode (black/yellow) for low-vision users",
            "  • WCAG 2.1 compliant semantic HTML",
            "  • Senior-citizen-friendly large fonts & simple flows",
            "",
            "🔒 SECURITY ARCHITECTURE",
            "  • 180-sec auto session purge — prevents identity theft at public kiosks",
            "  • SHA-256 Integrity Ledger — tamper-proof transaction & complaint hashes",
            "  • Anomaly Detection — flags brute-force & rapid-fire submissions",
            "  • Stateless frontend — zero sensitive data persists locally",
            "  • DPDP Act compliant — privacy-first design",
            "  • Rate limiting design for fraud prevention",
            "",
            "📋 COMPLIANCE",
            "  • Government UI guidelines followed",
            "  • WCAG 2.1 AA accessibility standard",
            "  • MeitY / Smart City Mission aligned",
        ], font_size=9.5, color=(30, 30, 30))

# ─────────────────────────────────────────────
# SLIDE 6 – Innovation & Future Scope
# ─────────────────────────────────────────────
slide6 = prs.slides[5]

add_multi_para_box(slide6, [
    (True,  "🚀 What Makes SUVIDHA Connect Unique"),
    (False, ""),
    (True,  "🤖 AI SUVIDHA Sahayak"),
    (False, "  NLP chat assistant handles Hindi & English queries,"),
    (False, "  navigates users directly to relevant service sections"),
    (False, ""),
    (True,  "🎤 Voice Commander"),
    (False, "  Fully hands-free kiosk operation via Web Speech API"),
    (False, "  Critical for hygiene & accessibility in public spaces"),
    (False, ""),
    (True,  "🏆 Gamification & Rewards"),
    (False, "  SUVIDHA Points for on-time payments & civic reporting"),
    (False, "  Redeemable for transport passes & municipal tax rebates"),
    (False, "  Drives digital adoption in low-literacy communities"),
    (False, ""),
    (True,  "🔗 Blockchain-Lite Integrity Ledger"),
    (False, "  SHA-256 hash for every transaction & complaint"),
    (False, "  Prevents administrative tampering of records"),
    (False, ""),
    (True,  "📱 PWA — Offline Capable"),
    (False, "  Installable on kiosk hardware, works without internet"),
    (False, "  Critical for rural & low-connectivity deployments"),
], left=311699, top=650000, width=4400000, height=4100000, font_size=10, color=(30, 30, 30))

add_multi_para_box(slide6, [
    (True,  "🔭 Future Roadmap"),
    (False, ""),
    (True,  "Phase 1 (0–6 months)"),
    (False, "  • Live API integration with Assam Gas & APDCL"),
    (False, "  • UPI / Bharat BillPay payment gateway"),
    (False, "  • PostgreSQL backend for persistent data"),
    (False, "  • SMS OTP via MSG91 / AWS SNS"),
    (False, ""),
    (True,  "Phase 2 (6–12 months)"),
    (False, "  • Expand to 5+ regional languages (Assamese, Bengali, Tamil)"),
    (False, "  • Aadhaar-based eKYC authentication"),
    (False, "  • DigiLocker integration for document fetch"),
    (False, "  • Predictive analytics for admin dashboard"),
    (False, ""),
    (True,  "Phase 3 (12–24 months)"),
    (False, "  • Pan-India rollout across Smart City kiosks"),
    (False, "  • IoT sensor integration (smart meters, water sensors)"),
    (False, "  • AI-powered grievance auto-routing to departments"),
    (False, "  • National Civic Services API federation"),
    (False, ""),
    (True,  "💡 Why We Will Win"),
    (False, "  Production-ready • Accessible • Secure • Scalable"),
    (False, "  Built for Bharat, not just metros"),
], left=4900000, top=650000, width=4000000, height=4100000, font_size=10, color=(30, 30, 30))

prs.save("SUVIDHA_Connect_Presentation.pptx")
print("✅ Saved: SUVIDHA_Connect_Presentation.pptx")
